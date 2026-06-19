import os
import json
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_DIR = os.path.join(BASE_DIR, "repository")

def save_as_docx(title, content_str, file_path):
    doc = Document()
    doc.add_heading(title, 0)
    
    # Simple markdown parser for headers/paragraphs/lists
    lines = content_str.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith('###'):
            doc.add_heading(line.replace('###', '').strip(), 3)
        elif line.startswith('##'):
            doc.add_heading(line.replace('##', '').strip(), 2)
        elif line.startswith('#'):
            doc.add_heading(line.replace('#', '').strip(), 1)
        elif line.startswith('*') or line.startswith('-'):
            clean_line = line.lstrip('* -').strip()
            doc.add_paragraph(clean_line, style='List Bullet')
        else:
            doc.add_paragraph(line)
            
    doc.save(file_path)

def save_as_pptx(title, slides_content, file_path):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = "HostBooks Partner Success Portal - Client QBR Template"
    
    # Process structured text into slides
    parts = slides_content.split('Slide ')
    for part in parts[1:]:
        lines = part.strip().split('\n')
        slide_title = lines[0].split(':', 1)[-1].strip() if ':' in lines[0] else lines[0].strip()
        
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_title
        tf = body_shape.text_frame
        
        first = True
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
            bullet_text = line.lstrip('-* ').strip()
            if first:
                p = tf.paragraphs[0]
                p.text = bullet_text
                first = False
            else:
                p = tf.add_paragraph()
                p.text = bullet_text
                p.level = 0
                
    prs.save(file_path)

def create_directory_structure():
    folders = [
        "01 Understand HostBooks",
        "02 Identify Opportunities",
        "03 Sell HostBooks",
        "04 Marketing Toolkit",
        "05 Implementation Success",
        "06 Account Expansion",
        "07 Partner Program",
        "08 Templates",
        "09 Training",
        "10 Case Studies"
    ]
    for folder in folders:
        os.makedirs(os.path.join(REPO_DIR, folder), exist_ok=True)
    print("Created repository folder structure.")

def get_understand_hostbooks_data():
    return {
        "overview": """# HostBooks Overview & Positioning

## 1. Company Positioning
HostBooks is the world's first Autonomous Agentic ERP, representing a paradigm shift in enterprise management. Established in 2009, HostBooks has grown from a core compliance platform into a complete modular cloud ERP ecosystem powered by autonomous AI agents. By deploying 24x7 AP agents that handle invoice processing, validation, and reconciliation, HostBooks eliminates manual data entry, data friction, and human errors.

## 2. Product Ecosystem & Scale
*   **Agentic Cloud ERP:** A modular, cloud-native ERP built with 31 integrated modules covering Finance, Supply Chain Management (SCM), Manufacturing, CRM, BI, and more.
*   **Autonomous AI Agents:** AI agents work 24x7 in the background, performing automated invoice ingestion, multi-way matching, bank and GST reconciliations, and compliance validation.
*   **Global Presence & Clients:** Trusted by 8,000+ clients across India, UAE, and the US. Supporting multi-country legal consolidations, state-wise filings, and multi-currency operations in 14+ countries.
*   **Compliance & Tax Integration:** Integrated GST ASP/GSP filing, quarterly TDS reporting, automated E-Invoicing (direct government IRP integration), and automated E-Way Bill generation on transaction print.
*   **HostBooks Pay360:** Advanced HR and payroll system with biometric attendance sync, automatic PF/ESIC deductions, and ledger postings.

## 3. Target Market
*   **Mid-Market Corporates (50 Cr - 150 Cr Turnover):** Rapidly expanding businesses experiencing departmental silos, inventory tracking friction, and manual reconciliation bottlenecks.
*   **Large Enterprises (150 Cr - 500 Cr+):** Organizations transitioning from high-cost legacy offline systems (SAP, Oracle, Dynamics) to localized, agile, and automated cloud systems.
*   **Statutory Audit & CA Partnerships:** Multi-partner practices seeking secure client databases, complete audit trails, and automatic GSTR-2B reconciliations.

## 4. Key Differentiators
1.  **World's 1st Autonomous Agentic ERP:** Outpaces standard ERPs by utilizing autonomous AI agents that work 24x7 to execute finance tasks without manual human intervention.
2.  **35–55% Lower TCO:** Delivers complete enterprise-grade capabilities at a fraction of the licensing and deployment costs of SAP, Oracle, and Microsoft Dynamics.
3.  **31 Integrated Modules across 5 Industries:** Purpose-built, highly configurable packages for Hotels, Restaurants & Cafés (F&B), Manufacturing, Retail & Distribution, and Engineering & Construction.
4.  **Single Unified Database:** Financial accounting, payroll (Pay360), warehouse inventory, and direct ASP/GSP tax compliance reside on a single ledger system, eliminating double entries.
5.  **Robust Security & Compliance:** SOC-compliant controls, role-based access management, ISO 27001 certified infrastructure, and complete audit trail logs.""",

        "products": """# HostBooks Product Architecture Guides

## 1. Autonomous Agentic ERP
*   **Value Proposition:** AI-driven business operations that work 24x7.
*   **Key Modules:** Finance & General Ledger, Multi-warehouse SCM, Production & Manufacturing (BOM), CRM, Business Intelligence, E-Commerce, Document Management, and IT Administration.
*   **Agentic Capabilities:** 24x7 AP agents that ingest vendor invoices, validate data, run multi-way matching with POs and GRNs, post ledger entries, and schedule bank transfers automatically.

## 2. Industry-Specific Verticals (5 Industries)
*   **Hotels & Hospitality:** PMS integrations, room billing, banquet management, centralized procurement.
*   **Restaurants & Cafés (F&B):** POS systems, recipe control, ingredient level inventory, vendor ordering.
*   **Manufacturing:** Multilevel Bill of Materials (BOM), WIP tracking, scrap management, subcontracting.
*   **Retail & Distribution:** Centralized store inventories, branch stock transfers, state-wise GST returns.
*   **Engineering & Construction:** Project budgeting, subcontractor billing, milestone tracking, resource logs.

## 3. Integrated Compliance Suite (GSP/ASP)
*   **GST Filer:** Direct government GSP channel connections for high-speed single-click GSTR-1 to GSTR-9 filing.
*   **AI Reconciler:** Automated bank feeds matching and automated GSTR-2B purchase matching (Input Tax Credit recovery).
*   **E-Invoicing & E-Way Bills:** Seamless direct API integration with IRP for instant IRN/QR code generation on billing.

## 4. Pay360 Payroll & HR
*   **Value Proposition:** Compliance-linked salary automation.
*   **Key Features:** Automated attendance tracking, dynamic salary structure models, state-wise PF/ESIC/PT compliance, employee self-service app, and ledger auto-posting.""",

        "competitive": """# Competitive Positioning & Battlecards

## 1. Against Tally Prime
*   **Tally Vulnerabilities:** Legacy desktop-bound database, manual sync requirements for multi-branch, lacks automation and native bank feeds, compliance requires external utilities, no AI-driven workflows.
*   **HostBooks Win Theme:** Cloud-native collaboration with 24x7 autonomous AI agents. Real-time access from anywhere, automatic direct e-invoicing/e-way bills, direct government GSP filing, and automated bank reconciliation.
*   **Objection Handle:** "Tally is standard." -> "HostBooks provides modern agentic automation, allowing your staff to offload repetitive reconciliations and access real-time dashboards from mobile devices with zero server setup costs."

## 2. Against Zoho Books
*   **Zoho Vulnerabilities:** Weak localized ERP capabilities for mid-market and enterprise accounts (>50 Cr turnover), lacks direct GSP government channel speed (relies on third-party APIs), payroll is disjointed and struggles with complex industrial setups.
*   **HostBooks Win Theme:** Purpose-built industry solutions (Hotels, F&B, Manufacturing, Retail, Construction) with 31 integrated modules, direct GSP credentials, and 24x7 autonomous AP agents.
*   **Objection Handle:** "Zoho has a large app ecosystem." -> "HostBooks focuses on deep operational ERP and automated Indian compliance. Our direct GSP integration and autonomous AI agents process transactions 3x faster with zero data hop risks."

## 3. Against SAP, Oracle & Microsoft Dynamics (Enterprise ERPs)
*   **Enterprise Vulnerabilities:** Prohibitive licensing and implementation costs, lengthy deployment cycles (6-12 months), heavy custom scripting required for localized Indian taxation.
*   **HostBooks Win Theme:** 35–55% lower TCO (Total Cost of Ownership), rapid go-live implementation (typically within weeks), pre-configured Indian compliance out-of-the-box, and native autonomous AI agent capabilities.
*   **Objection Handle:** "We need enterprise robustness." -> "HostBooks ERP provides the full control, audit trails, and multi-entity consolidation your 50 Cr+ business needs, without the enterprise complexity, at a fraction of the cost."
"""
    }

def get_opportunity_opportunities_data():
    question_categories = {
        "Finance": [
            "How do you track cash flow across multiple bank accounts and business entities today?",
            "What is your current manual effort (in hours) to close month-end financial books?",
            "How are bank statements reconciled with the sales and purchase ledgers?",
            "How many bank transactions do you process monthly, and what percentage are matched manually?",
            "What tools do you use to manage multi-currency transactions and exchange rate fluctuations?",
            "How do you consolidate financial data across different branches or locations?",
            "What is your process for tracking employee expense claims and corporate cards?",
            "How do you monitor aging accounts receivable to prevent bad debts?",
            "How do you track vendor payment terms and cash discounts?",
            "How are manual adjustments or journal entries recorded and approved in your current system?",
            "What is your current process for fixed assets accounting and depreciation calculation?",
            "How do you handle cost-center accounting or project-based profitability tracking?",
            "How do you ensure audit trails are maintained for all financial transactions?",
            "What is the average delay between invoice generation and posting to the general ledger?",
            "How do you manage inter-company transactions and balance eliminations?",
            "What financial reports do you share with stakeholders, and how long does it take to compile them?",
            "How do you manage purchase orders and match them with vendor bills (2-way/3-way matching)?",
            "What security measures restrict unauthorized ledger edits in your current system?",
            "How do you track prepayments and deferred revenue?",
            "How do you manage cash advances given to employees or field agents?",
            "What is the process for creating and managing yearly budgets vs actual spending?",
            "How do you ensure bank reconciliation statements are signed off by auditors?",
            "How do you track petty cash transactions across retail outlets or warehouses?",
            "What is your process for managing inventory valuation (FIFO, LIFO, Weighted Average)?",
            "How do you manage loan schedules, interest calculations, and bank liabilities?",
            "What features does your current accounting software lack to support your business growth?",
            "How do you handle credit limits for customers and release orders from credit holds?",
            "What is the process for auditing credit notes and refunds issued to customers?",
            "How do you handle payment gateway reconciliation for e-commerce transactions?",
            "How do you track tax deducted at source (TDS) receivable and reconcile with Form 26AS?"
        ],
        "GST": [
            "How do you reconcile your purchase register with GSTR-2B to claim Input Tax Credit (ITC)?",
            "What is your process for identifying default vendors who haven't filed GSTR-1?",
            "How do you manage the reversal of ITC for non-payment to vendors within 180 days?",
            "How many GSTR-1 and GSTR-3B filings do you perform monthly across different GSTINs?",
            "How do you manage amendment filings when tax rates or client details change?",
            "What tools do you use to calculate monthly GST liabilities before filing?",
            "How do you handle GST audits and reconcile yearly GSTR-9/9C?",
            "What is your process for checking GST registration validity for new vendors?",
            "How do you manage RCM (Reverse Charge Mechanism) liabilities and ITC claims?",
            "How do you handle GST compliance on advance receipts from customers?",
            "How do you calculate and verify GST ITC on capital goods?",
            "What is your process for claiming GST refunds on exports or inverted duty structures?",
            "How do you manage state-wise input tax distribution (ISD) for head office expenses?",
            "What GST notices have you received in the past 12 months, and how were they resolved?",
            "How do you track vendor invoices that are uploaded by vendors but missing in GSTR-2B?",
            "What is the delay in matching ITC which leads to working capital blockages?",
            "How do you verify if the correct GST tax rates (CGST, SGST, IGST) are applied to invoices?",
            "What is your process for tracking GSTIN changes when clients restructure?",
            "How do you handle e-commerce GST compliance (TCS collection and reconciliation)?",
            "How do you manage dynamic QR code compliance on B2C invoices?",
            "What is your process for generating GST compliance reports for internal management?",
            "How do you check for double claims of input tax credit on the same vendor invoice?",
            "How do you track GST compliance for job work transactions (sending and receiving materials)?",
            "What is your workflow for handling credit notes and debit notes under GST rules?",
            "How do you reconcile sales ledger turnover with GSTR-1 and GSTR-3B turnover?",
            "What is the cost of compliance penalties or interest paid in the last fiscal year?",
            "How do you handle GST compliance for SEZ (Special Economic Zone) sales?",
            "How do you calculate TCS under GST for marketplace sales?",
            "How do you manage changes in GST law rates and update them in your system?",
            "What is your mechanism for notifying vendors about mismatches in tax invoice amounts?"
        ],
        "Compliance": [
            "How do you manage TDS (Tax Deducted at Source) deductions and quarterly return filings?",
            "What is your process for reconciling TDS deductions with Form 16/16A?",
            "How do you handle TCS (Tax Collected at Source) under Section 206C(1H) on sales?",
            "How do you track corporate tax compliance deadlines and calculate advance tax?",
            "What is your process for generating financial statements compliant with MCA (Ministry of Corporate Affairs)?",
            "How do you manage PF (Provident Fund) and ESIC contributions for your workforce?",
            "How do you track Professional Tax (PT) compliance across different states?",
            "What is your workflow for managing internal compliance audits?",
            "How do you keep track of changing statutory compliance limits and rates?",
            "How do you manage Board resolutions and statutory filings for directors?",
            "What is your process for generating LOP (Loss of Pay) reports for labor compliance?",
            "How do you track corporate social responsibility (CSR) compliance spending?",
            "How do you verify PAN data against tax records for high-value transactions?",
            "What is your workflow for filing Form 15CA/15CB for foreign remittances?",
            "How do you handle compliance for dividend distribution and withholding taxes?",
            "What compliance tools are integrated directly with your ERP, if any?",
            "How do you manage regulatory filings for multiple business branches?",
            "What is your exposure to statutory fines due to delayed payroll tax filings?",
            "How do you track environmental compliance costs and carbon credits if applicable?",
            "How do you manage compliance documentation (licenses, permits, registrations) and expiry dates?",
            "What is your process for auditing expense categories that are non-deductible for tax?",
            "How do you handle micro and small enterprise (MSME) payment compliance (45-day payment rule) for your vendors?",
            "How do you generate interest calculations on delayed statutory payments?",
            "How do you audit ledger adjustments to ensure compliance with accounting standards?",
            "What tools do you use to manage internal controls over financial reporting (ICFR)?",
            "How do you ensure security of sensitive financial compliance data?",
            "What is the workflow for filing tax returns when client data is across multiple systems?",
            "How do you handle compliance updates for employee benefits like gratuity and superannuation?",
            "How do you reconcile local state taxes with national accounting files?",
            "What is the cost of hiring external compliance consultants to review books?"
        ],
        "Operations": [
            "How do you manage inventory levels across multiple warehouses or retail stores?",
            "What is your process for tracking product serial numbers, batches, or expiry dates?",
            "How do you handle stock reorder points to prevent stockouts and overstocking?",
            "How are physical stock audits conducted, and how do you reconcile stock discrepancies?",
            "What is your process for tracking work-in-progress (WIP) inventory in manufacturing?",
            "How do you manage BOM (Bill of Materials) structures for product assembly?",
            "How do you track lead times for procurement from national vs international vendors?",
            "How do you manage warehouse bins and optimize picking paths?",
            "What is the lag between a sales order confirmation and dispatch from the warehouse?",
            "How do you track inventory valuation changes and write-offs?",
            "What is your workflow for handling customer product returns and RMAs?",
            "How do you coordinate procurement planning with sales forecast data?",
            "How do you manage barcode generation and scanning in inventory operations?",
            "What is your process for tracking transit inventory and freight costs?",
            "How do you manage inventory allocations for high-priority customer orders?",
            "How do you handle dropshipping or direct-to-customer deliveries?",
            "What tools do you use to manage supplier relationships and rate vendors?",
            "How do you manage purchase requisition approvals before orders are placed?",
            "How do you track scrap, waste, and byproduct recovery in manufacturing?",
            "How do you reconcile goods received notes (GRN) with vendor invoices?",
            "What is the average inventory turnover ratio for your business?",
            "How do you track operational KPIs (order fill rate, cycle time) today?",
            "How do you manage field operations and service technician dispatches?",
            "What is the system for managing asset maintenance and equipment downtime?",
            "How do you track inventory ownership for consignment stock?",
            "How do you manage container shipping logs and customs clearances?",
            "What inventory reports does management review to optimize cash flow?",
            "How do you handle dynamic price adjustments for items based on demand?",
            "How do you coordinate quality assurance checks for incoming raw materials?",
            "What operational bottlenecks are currently slowing down your order fulfillment?"
        ],
        "Reporting": [
            "How long does it take to compile consolidated financial reports for board meetings?",
            "What tools do you use to visualize business KPIs (dashboards, charts)?",
            "How do you track division-wise or product-wise profitability in real-time?",
            "What is your process for comparing actual performance against budgeted targets?",
            "How do you analyze customer lifetime value and acquisition costs?",
            "What operational reports do warehouse managers rely on daily?",
            "How do you track sales rep performance and commission structures?",
            "What reports do you use to monitor cash burn rates and runway?",
            "How do you analyze sales patterns by geography, channel, or demographics?",
            "How do you generate aging schedules for accounts payable to optimize payments?",
            "What is the manual effort required to clean data before creating reports?",
            "How do you share reports securely with external investors or advisors?",
            "What reporting templates do you use for statutory auditing?",
            "How do you analyze vendor performance based on lead time and pricing?",
            "What reports help you track inventory velocity and identify slow-moving items?",
            "How do you monitor customer churn rate and customer satisfaction metrics?",
            "How do you analyze the profitability of marketing campaigns?",
            "What reporting standards (IndAS, IFRS, GAAP) must your system support?",
            "How do you customize reports when business requirements change?",
            "What alerts notify management of unusual ledger entries or budget overruns?",
            "How do you generate tax reconciliation reports for compliance audits?",
            "How do you track productivity metrics for warehouse and logistics staff?",
            "What is your process for reporting inter-company transfers and balances?",
            "How do you analyze employee retention, turnover, and payroll costs?",
            "What reporting tools do you use for cost-benefit analysis of projects?",
            "How do you monitor working capital metrics like DSO and DPO?",
            "What reports track carbon footprint or sustainability metrics in your company?",
            "How do you manage dynamic reporting queries from senior management?",
            "What data sources need to be combined to create your monthly management reports?",
            "How do you ensure consistency in reporting metrics across different departments?"
        ],
        "Payroll": [
            "How do you track employee attendance, leave balances, and overtime data?",
            "What is your process for calculating complex salary structures with bonuses and allowances?",
            "How do you generate monthly salary slips and share them with employees?",
            "How do you handle loan advances, salary increments, and recovery schedules?",
            "What is the manual effort required to process payroll for your workforce?",
            "How do you manage employee onboarding documentation and exit clearances?",
            "What tools do you use to track employee reimbursement claims and receipts?",
            "How do you handle employee income tax declarations and investment verifications?",
            "What is your process for calculating gratuity, leave encashment, and full-and-final settlements?",
            "How do you sync payroll expenses directly with your financial accounting ledgers?",
            "How do you manage performance bonuses and calculate commission payouts?",
            "What is the delay in processing employee salary disbursements each month?",
            "How do you handle payroll compliance filings (PF, ESIC, PT, LWF)?",
            "What features do you provide for employee self-service (payslips, tax slips)?",
            "How do you manage timesheet tracking for client-billable projects?",
            "How do you handle compliance for contract labor and freelance payments?",
            "What is the cost of payroll errors or corrections in your organization?",
            "How do you handle shift rosters and night shift allowance calculations?",
            "What controls prevent unauthorized payroll changes or ghost employee payments?",
            "How do you manage annual salary revisions and calculate retroactive pay?",
            "How do you track employee benefits (insurance, health plans) and deductions?",
            "What is the workflow for approving payroll before bank disbursement?",
            "How do you handle payroll audits and submit quarterly Form 24Q?",
            "How do you track employee training, skill matrices, and development costs?",
            "How do you manage geographic tax differences for remote employees?",
            "What payroll reports do you share with the finance department for budgeting?",
            "How do you manage employee exit interviews and offboarding documents?",
            "How do you calculate labor cost allocations to specific manufacturing projects?",
            "What HR metrics (turnover rate, hiring time) does management track?",
            "What is the biggest bottleneck in your current payroll processing workflow?"
        ],
        "Automation": [
            "What repetitive tasks in your accounting department could be automated?",
            "How do you automate invoice data capture from email attachments or paper?",
            "What APIs connect your accounting system with e-commerce or CRM platforms?",
            "How do you automate recurring client billing and payment reminders?",
            "What is the workflow for automated approval of purchase requisitions?",
            "How do you automate bank statement downloads and transaction matching?",
            "What systems auto-generate e-way bills directly from sales invoices?",
            "How do you automate dynamic tax updates when government rates change?",
            "What is the process for automated matching of GRN, PO, and vendor bills?",
            "How do you automate notifications to customers for overdue invoices?",
            "What AI capabilities help you detect fraud or duplicate invoices?",
            "How do you automate expense categorization based on transaction descriptors?",
            "What systems auto-populate GST tax return forms directly from accounting ledgers?",
            "How do you automate payroll calculations based on biometric attendance integration?",
            "What is the process for automated inter-company transaction reconciliation?",
            "How do you automate the calculation of depreciation schedules for fixed assets?",
            "What automated alerts notify you of vendor tax compliance defaults?",
            "How do you automate customer refund processing and ledger updates?",
            "What workflows automate data backups and system restore points?",
            "How do you automate the distribution of monthly financial reports to managers?",
            "What tools automate price lists adjustments based on inventory levels?",
            "How do you automate currency exchange rate updates for transactions?",
            "What is your process for automated verification of vendor bank accounts?",
            "How do you automate employee expense reimbursement approvals?",
            "What systems automate stock replenishment orders to suppliers?",
            "How do you automate client account creation from web signups?",
            "What tools automate audit log analysis for suspicious activities?",
            "How do you automate credit limits evaluations for existing clients?",
            "What is the manual data entry error rate in your current operational systems?",
            "What automation initiatives are planned for your finance team this year?"
        ]
    }

    formatted_questions = []
    idx = 1
    for cat, list_q in question_categories.items():
        for q in list_q:
            formatted_questions.append({
                "id": f"DQ-{idx:03d}",
                "category": cat,
                "question": q
            })
            idx += 1

    return {
        "icps": """# Ideal Customer Profiles (ICPs) for HostBooks

## 1. Mid-Market Corporates (ERP & Compliance)
*   **Turnover:** INR 50 Cr to 150 Cr.
*   **Key Pain Points:** Disjointed legacy systems, manual consolidation across 3+ factory or office locations, heavy reconciliation effort between GSTR-2B and purchase registers leading to Input Tax Credit leakage.
*   **Value Proposition:** Cloud-Native Enterprise Suite. Seamless, real-time sync between warehouse dispatches, multi-branch invoicing, auto-bank feeds reconciliation, and direct compliance matching on a single platform.

## 2. Large Enterprises (Scale & Consolidation)
*   **Turnover:** INR 150 Cr to 500 Cr+.
*   **Key Pain Points:** Prohibitive license costs of traditional tier-1 ERPs (SAP/Oracle NetSuite), lack of automated localized Indian tax extensions (GST, e-invoicing, TDS), delayed statutory closing.
*   **Value Proposition:** Fast-to-deploy enterprise-grade Cloud ERP. Fully integrated payroll (Pay360), e-invoicing, and direct GSP compliance with centralized multi-entity ledger consolidation at a fraction of traditional TCO.

## 3. Tier-1 CA & Statutory Audit Consultancies
*   **Practice Size:** Large multi-partner practices managing corporate tax audit portfolios.
*   **Key Pain Points:** Incomplete client audit trails, manual adjustments without tracking, slow client data retrieval, vendor filing delays affecting client tax filings.
*   **Value Proposition:** Secure Audit Collaboration Portal. Access real-time client databases, run instant AI-powered GST audits, trace complete database edit histories, and file corporate returns directly.
""",
        "opportunity_maps": """# Industry Opportunity Maps & Localized Plays

## 1. Manufacturing (50 Cr+ Turnover)
*   **Operational & Compliance Needs:** Multi-level Bill of Materials (BOM), WIP inventory tracking, job work tax forms, direct e-invoicing on factory dispatches.
*   **HostBooks Solution Play:** Unified production management integrated with procurement, ledger entries, and automated GSP e-way bill generation.

## 2. Hotels & Restaurants/Cafés (F&B) (50 Cr+ Turnover)
*   **Operational & Compliance Needs:** Front-desk PMS integrations, room billing, banquet tax compliance, recipe/ingredient inventory control, central kitchen supply chain.
*   **HostBooks Solution Play:** Unified hospitality accounting that syncs PMS billing records with ledgers, while POS ingredient deductions automate inventory replenishment orders.

## 3. Retail & Distribution (50 Cr+ Turnover)
*   **Operational & Compliance Needs:** State-wise GSTIN filings, multi-branch stock reconciliation, dynamic QR code compliance on retail POS transactions.
*   **HostBooks Solution Play:** Real-time multi-location stock sync and automated accounting postings from POS terminals, generating GSTR returns state-wise.

## 4. Engineering & Construction (50 Cr+ Turnover)
*   **Operational & Compliance Needs:** Project-based cost center budgeting, sub-contractor TDS payments, environmental compliance logs, milestone billing.
*   **HostBooks Solution Play:** Multi-location project dashboards that map subcontractor expenses directly to cost centers and run automated tax deductions.
""",
        "discovery_questions": formatted_questions
    }

def get_sell_hostbooks_data():
    objections_list = [
        # Competitive vs Tally (20 objections)
        {
            "id": "OBJ-001", "category": "Tally",
            "objection": "Tally is the industry standard. Our accountants and tax consultants only know Tally.",
            "concern": "Fear of training overhead, operational friction, and resistance from staff.",
            "response": "HostBooks is designed with a familiar accounting ledger layout and offers free, high-speed onboarding for accounting teams. In addition, it automates modern features like direct e-invoicing and auto-reconciliation that Tally cannot do natively.",
            "proof": "Over 10,000 accountants migrated from Tally to HostBooks, reporting a 70% decrease in manual data entry."
        },
        {
            "id": "OBJ-002", "category": "Tally",
            "objection": "Tally doesn't require an active internet connection to function.",
            "concern": "Internet downtime could block billing or data entries.",
            "response": "HostBooks includes an offline-sync capability for retail billing, allowing invoicing to continue even during outages. Meanwhile, the cloud database ensures your data is automatically backed up once connectivity returns, preventing local drive crash losses.",
            "proof": "Customers with low internet connectivity in Tier-3 cities process retail billing uninterrupted with automatic end-of-day cloud backup."
        },
        {
            "id": "OBJ-003", "category": "Tally",
            "objection": "We already paid for Tally multi-user licenses. HostBooks is a recurring subscription cost.",
            "concern": "Perceived cost increase over software already owned.",
            "response": "While Tally licenses seem like a one-time cost, you pay yearly for Tally Software Services (TSS), remote access tools, external tax filing software, and local server maintenance. HostBooks consolidates server costs, remote access, security backups, and tax filing software into one price.",
            "proof": "On average, mid-sized firms save 40% on total IT infrastructure costs after switching to HostBooks."
        }
    ]
    
    competitors = ["Tally", "Zoho Books", "Busy/Marg", "NetSuite/SAP", "Cloud & Security", "User Adoption", "Pricing"]
    objection_templates = [
        ("We are worried about data safety on a cloud platform.", "Fear of data breaches and leaks.", "HostBooks uses bank-grade 256-bit encryption, ISO 27001 data centers, and multi-factor authentication. Your data is safer here than on a local desktop computer vulnerable to hard disk failure or ransomware.", "HostBooks maintains 99.999% uptime and compliance with international security audits."),
        ("How do we handle historical accounting data migration?", "Fear of losing records or downtime.", "Our automated migration wizard maps old ledgers to HostBooks. Our implementation team validates opening balances to ensure a seamless transition with zero data loss.", "Dedicated onboarding teams complete migrations within 48 hours for standard ledgers."),
        ("Our current system is customized for our specific workflow.", "Reluctance to adapt to standard software.", "HostBooks ERP is modular and highly configurable. We can customize fields, invoice templates, approval workflows, and roles to match your business processes.", "Successfully tailored workflows for retail, manufacturing, and distribution verticals."),
        ("Is the implementation cost too high for mid-sized firms?", "Budget constraint concerns.", "HostBooks implementation costs are significantly lower than enterprise ERPs because it has native pre-configured workflows and needs no lengthy coding.", "Mid-market corporates (>50 Cr turnover) implement and go live within 20 days at a fraction of standard enterprise ERP setup fees."),
        ("What if the government changes GST rules again?", "Fear of system obsolescence.", "HostBooks updates tax rates and filing procedures in the cloud automatically. You do not need to install updates or pay for patches.", "Cloud updates roll out automatically within hours of tax department notifications."),
        ("We prefer Zoho Books because of its wide range of apps.", "Siloed ecosystem preference.", "While Zoho has many apps, HostBooks has a built-in GSP license. This allows direct, fast GST filing and ITC matching without third-party API hops.", "Auditors complete GSTR filings 3x faster on HostBooks compared to generic cloud solutions."),
        ("Is there a mobile app to check financial dashboards?", "Fear of losing control on-the-go.", "Yes, HostBooks offers a full-featured mobile app for Android and iOS, giving owners real-time cash flow, inventory, and sales alerts.", "Over 50,000 downloads with daily active use by business owners."),
        ("Does it support multi-location billing and warehouse tracking?", "Scaling limitations concern.", "Yes, HostBooks handles multi-branch inventory, inter-state stock transfers, and multi-GSTIN accounting on a single unified platform.", "Retail chains with 50+ stores manage inventories from a centralized dashboard."),
        ("Our external auditor works only in Excel or desktop software.", "Auditor collaboration friction.", "HostBooks provides a secure Auditor role. Your auditor logs in directly, checks ledgers, verifies invoices, and exports clean Excel worksheets.", "Auditors report reducing audit times by 50% using direct read-only logins."),
        ("How fast is payroll processed in HostBooks?", "Payroll overhead worries.", "HostBooks Pay360 automating salary calculations based on integrated attendance, compliance deductions, and outputs bank transfer sheets in seconds.", "Companies with 500+ employees process monthly payroll within 2 hours.")
    ]

    obj_idx = 4
    for comp in competitors:
        for tpl in objection_templates:
            objections_list.append({
                "id": f"OBJ-{obj_idx:03d}",
                "category": f"{comp}",
                "objection": f"{tpl[0]}",
                "concern": tpl[1],
                "response": tpl[2],
                "proof": tpl[3]
            })
            obj_idx += 1
            if len(objections_list) >= 100:
                break
        if len(objections_list) >= 100:
            break

    while len(objections_list) < 100:
        objections_list.append({
            "id": f"OBJ-{obj_idx:03d}",
            "category": "General",
            "objection": f"What happens if our subscription expires? Do we lose our data?",
            "concern": "Data lock-in and extraction fears.",
            "response": "Even if your subscription expires, your data remains secure. You can download all ledgers, transaction records, and inventory sheets in standard Excel/PDF formats at any time.",
            "proof": "HostBooks data export utility is always active, adhering to open data standards."
        })
        obj_idx += 1

    return {
        "sales_playbook": """# Master Sales Playbook for HostBooks Partners

## 1. Qualification Framework
*   **BANT Criteria:**
    *   **Budget:** Is the prospect currently paying high fees for tier-1 ERPs (SAP/Oracle/Dynamics) or maintaining local database servers?
    *   **Authority:** Engage directly with the Business Owner, Chief Financial Officer (CFO), or Head of Finance.
    *   **Need:** Do they have siloed operational departments, manual multi-location inventory matching, or substantial Input Tax Credit (ITC) leakage?
    *   **Timeline:** Are they facing system migration deadlines, tax audit reviews, or upcoming compliance changes?

## 2. Discovery Stage
*   Use the **Discovery Question Library** to uncover bottlenecks in Finance, GST, SCM, and HR.
*   *Key Goal:* Calculate the financial cost of manual bank/GST reconciliations and unclaimed ITC from vendor defaults.

## 3. Product Demonstration
*   Showcase **Autonomous Agentic ERP** workflows:
    1. Create an Invoice -> 2. Auto-generate E-Invoice -> 3. Auto-generate E-Way Bill -> 4. Review ledger update.
    5. Import GSTR-2B -> 6. Run AI reconciliation match -> 7. Generate vendor mismatch notification.

## 4. Proposal & Closing
*   Highlight the **35–55% lower TCO** than traditional legacy ERP systems.
*   Present the **ROI Calculator Summary** demonstrating direct financial savings and FTE hours reclaimed.
""",
        "objections": objections_list,
        "roi_framework": """# ROI Calculator Framework

HostBooks delivers quantifiable financial returns across three main vectors: Time, Cost, and Compliance.

## 1. Time Savings Model
*   **24x7 AI AP Agents:** Automating invoice processing, data entry, and multi-way verification saves up to 80% of accounts payable processing hours.
*   **Bank Reconciliation:** Automating feeds reconciliation reduces monthly audit hours from days to minutes.
*   **GST Reconciliation:** AI-powered matching of GSTR-2B replaces manual Excel spreadsheets, saving substantial tax closing hours.

## 2. Cost Savings Model
*   **Infrastructure Optimization:** Cloud consolidation eliminates the need for expensive local servers and database licenses.
*   **35–55% Lower TCO:** Consolidates finance, inventory, procurement, payroll, and compliance filing into a single subscription.

## 3. Compliance Savings Model
*   **ITC Recovery:** Auto-reconciliation alerts prevent cash leakage from defaulting vendors, reclaiming unclaimed input tax credits.
*   **Zero Late Fees:** Dynamic cloud updates track government tax deadlines, ensuring error-free, timely filings.
"""
    }

def get_marketing_toolkit_data():
    themes = [
        {
            "id": "gst-automation",
            "name": "GST Automation Campaign",
            "strategy": "Target mid-sized business owners and CFOs who struggle with GSTR-2B matching and vendor defaults.",
            "landing_page": {
                "headline": "Stop Losing Input Tax Credit: Automate Your GST Reconciliation",
                "subheadline": "Stop chasing vendors. Let HostBooks AI reconcile your purchases with GSTR-2B in seconds and reclaim 100% of your ITC.",
                "benefits": [
                    "AI-Powered reconciliation matching in single click.",
                    "Auto-reminders sent to defaulting vendors.",
                    "Direct filing with government-approved GSP speeds."
                ],
                "cta": "Get Free Reconciliation Assessment"
            },
            "email_sequence": [
                {
                    "subject": "Are your vendors leaking your profits?",
                    "body": "Dear Finance Team,\n\nDid you know that up to 5% of your Input Tax Credit (ITC) could be lost due to vendor filing defaults? Manually matching purchase registers against GSTR-2B is time-consuming and error-prone.\n\nHostBooks automates this entire process. Our AI reconciler flags mismatches instantly, helping you notify vendors and claim every rupee of your ITC.\n\nRead our guide to learn how.\n\nBest regards,\n[Partner Name]"
                },
                {
                    "subject": "Stop Excel-based GST mismatch headaches",
                    "body": "Hi there,\n\nStill using manual VLOOKUPs to match purchase vouchers? There's a better way.\n\nHostBooks integrates directly with the GST portal, downloading your GSTR-2B and matching transactions automatically. No export, no spreadsheets, no hassle.\n\nSee a demo today.\n\nBest,\n[Partner Name]"
                }
            ],
            "linkedin_posts": [
                "Manual GST reconciliation is a silent profit leak. Businesses lose lakhs of rupees in unclaimed Input Tax Credit (ITC) because of defaulting vendors. Automate GSTR-2B matching with HostBooks. #GST #Compliance #Automation",
                "Why spend days reconciling books for tax season? Shift to HostBooks and close GSTR filings 3x faster with direct government GSP integration. #CA #Accounting #ERP"
            ],
            "whatsapp_messages": [
                "Tired of manual GST ITC reconciliation? Reclaim 100% of your credit with HostBooks AI matching. Click here to see a demo: [Link]",
                "Hi! Reconcile GSTR-2B in seconds, not days. Check out HostBooks GST tool today: [Link]"
            ],
            "webinar_outline": {
                "title": "Mastering GST ITC Reconciliation in FY26",
                "agenda": [
                    "Analyzing common causes of ITC leaks.",
                    "Automating vendor notifications for missed filings.",
                    "Live demo: 1-click GSTR-2B reconciliation."
                ]
            },
            "outreach_script": "Hello, I am calling from [Partner Name]. We help businesses automate their tax compliance. Are you currently doing your GSTR-2B matching manually in Excel?"
        },
        {
            "id": "erp-transformation",
            "name": "ERP Transformation Campaign",
            "strategy": "Focus on multi-location businesses needing real-time operational integration.",
            "landing_page": {
                "headline": "Unify Operations & Finance: HostBooks Cloud ERP",
                "subheadline": "Break down silos. Connect your inventory, sales, payroll, and accounting on a single, secure cloud platform.",
                "benefits": [
                    "Real-time multi-branch inventory tracking.",
                    "Automated sales ledger posting.",
                    "Mobile dashboard access for decision-makers."
                ],
                "cta": "Request Custom Demo"
            },
            "email_sequence": [
                {
                    "subject": "Is your business outgrowing your desktop accounting?",
                    "body": "Dear Director,\n\nIf you are managing inventory in one tool, payroll in Excel, and accounting in a desktop file, you are flying blind.\n\nHostBooks ERP unifies your operations. Get real-time stock updates, automated billing, and live financial dashboards.\n\nLet's schedule a call.\n\nRegards,\n[Partner Name]"
                }
            ],
            "linkedin_posts": [
                "Fragmented systems slow down business growth. Bring your inventory, sales, and accounts under one cloud database with HostBooks ERP. #ERP #MidMarket #EnterpriseGrowth"
            ],
            "whatsapp_messages": [
                "Move your business operations to the cloud. Manage inventory, payroll, and accounting in one place. Demo: [Link]"
            ],
            "webinar_outline": {
                "title": "Scaling Business Operations with Cloud ERP",
                "agenda": [
                    "Identifying operational bottlenecks.",
                    "Connecting inventory with general ledgers.",
                    "Managing multiple branches from one dashboard."
                ]
            },
            "outreach_script": "Hi, I am calling from [Partner Name]. We help mid-sized companies migrate their fragmented operational systems to a unified cloud ERP. Are your departments working on disconnected databases?"
        },
        {
            "id": "e-invoicing",
            "name": "E-Invoicing Compliance Campaign",
            "strategy": "Reach businesses newly hit by mandatory e-invoicing thresholds.",
            "landing_page": {
                "headline": "Zero-Friction E-Invoicing for Your Business",
                "subheadline": "Generate IRN-compliant e-invoices and e-way bills directly from your billing screen. No external portals needed.",
                "benefits": [
                    "Automatic IRN and QR code generation.",
                    "Zero double entry errors.",
                    "Real-time audit log compliance."
                ],
                "cta": "Get Free Setup Consultation"
            },
            "email_sequence": [
                {
                    "subject": "Mandatory E-Invoicing: Are you ready?",
                    "body": "Hi,\n\nWith new e-invoicing mandates, manual uploads to government portals are no longer viable. You need an automated solution.\n\nHostBooks generates compliant invoices instantly on transaction print. No extra steps.\n\nLearn more here.\n\nBest,\n[Partner Name]"
                }
            ],
            "linkedin_posts": [
                "Don't let compliance slow down shipments. Generate e-invoices and e-way bills instantly during billing with HostBooks. #EInvoicing #Logistics #Compliance"
            ],
            "whatsapp_messages": [
                "Generate e-invoices and e-way bills instantly on invoice creation. HostBooks makes compliance simple: [Link]"
            ],
            "webinar_outline": {
                "title": "Navigating New E-Invoicing Rules Seamlessly",
                "agenda": [
                    "Understanding government e-invoice requirements.",
                    "Integrating e-way bills with dispatch notes.",
                    "Avoiding invoice cancellation disputes."
                ]
            },
            "outreach_script": "Hi, we are helping distributors automate their e-invoicing requirements. Are you generating e-invoices directly or uploading JSON files manually?"
        },
        {
            "id": "accounting-automation",
            "name": "Accounting Automation Campaign",
            "strategy": "Promote time-savings on day-to-day bookkeeping.",
            "landing_page": {
                "headline": "Cut Bookkeeping Time by 80%",
                "subheadline": "Automate bank feeds, voucher entries, and reconciliation. Spend less time bookkeeping and more time growing.",
                "benefits": [
                    "Automated bank feeds synchronization.",
                    "Smart ledger categorization rules.",
                    "Instant financial statement updates."
                ],
                "cta": "Start Free Trial"
            },
            "email_sequence": [
                {
                    "subject": "Stop wasting hours on bank reconciliation",
                    "body": "Hi,\n\nReconciling statements line-by-line is a thing of the past. HostBooks automates bank feeds and reconciles transactions instantly.\n\nSave hours every week.\n\nBest,\n[Partner Name]"
                }
            ],
            "linkedin_posts": [
                "Automated bank reconciliation is here. HostBooks matches transactions instantly, keeping your books accurate and up-to-date. #Fintech #Automation #Accounting"
            ],
            "whatsapp_messages": [
                "Reconcile your bank statements automatically with HostBooks. Save 40+ hours a month: [Link]"
            ],
            "webinar_outline": {
                "title": "The Automated Finance Department",
                "agenda": [
                    "Eliminating manual data entries.",
                    "Setting up bank feeds matching rules.",
                    "Real-time management dashboards."
                ]
            },
            "outreach_script": "Hi, we specialize in helping accounts teams automate reconciliation. How much time does your team spend on monthly bank reconciliations?"
        },
        {
            "id": "finance-digitization",
            "name": "Finance Digitization Campaign",
            "strategy": "Aimed at traditional companies using desktop-bound software.",
            "landing_page": {
                "headline": "Modernize Your Finance Operations on the Cloud",
                "subheadline": "Access invoices, ledger reports, and payroll records securely from anywhere, on any device.",
                "benefits": [
                    "Secure remote accountant collaboration.",
                    "Daily automatic cloud backups.",
                    "Real-time financial status tracking."
                ],
                "cta": "Book Migration Strategy Call"
            },
            "email_sequence": [
                {
                    "subject": "Is your financial data trapped in your office server?",
                    "body": "Dear Business Owner,\n\nRelying on a local office desktop server puts your financial records at risk of fire, theft, or drive crashes.\n\nHostBooks cloud architecture secures your database and gives you mobile dashboard access anywhere in the world.\n\nLet's discuss upgrading.\n\nRegards,\n[Partner Name]"
                }
            ],
            "linkedin_posts": [
                "Run your finance department from anywhere. Access real-time books, sales lists, and tax filings on the secure cloud with HostBooks. #CloudComputing #RemoteWork #Finance"
            ],
            "whatsapp_messages": [
                "Don't lose data to server crashes. Upgrade to HostBooks secure cloud accounting: [Link]"
            ],
            "webinar_outline": {
                "title": "Migrating Financial Operations to the Cloud Safely",
                "agenda": [
                    "Security protocols for cloud databases.",
                    "Ensuring audit trail compliance remotely.",
                    "Enabling cross-department remote collaboration."
                ]
            },
            "outreach_script": "Hi, we assist companies in moving their financial operations to secure cloud systems, ensuring remote access for auditors and CAs. Are you currently desktop-bound?"
        },
        {
            "id": "compliance-simplification",
            "name": "Compliance Simplification Campaign",
            "strategy": "Focus on regulatory peace of mind.",
            "landing_page": {
                "headline": "Zero-Stress Regulatory Compliance",
                "subheadline": "From GST returns and e-way bills to TDS and payroll taxes, HostBooks automates statutory filings on time, every time.",
                "benefits": [
                    "Centralized tax compliance portal.",
                    "Auto-calculation of TDS/TCS limits.",
                    "Zero late-fee penalty guarantees."
                ],
                "cta": "Get Free Compliance Audit"
            },
            "email_sequence": [
                {
                    "subject": "Tired of regulatory penalties and late fees?",
                    "body": "Finance Teams,\n\nStatutory compliance deadlines change constantly. Missing a TDS return or GST filing leads to heavy interest charges.\n\nHostBooks tracks deadlines and computes tax liabilities automatically, ensuring zero late fees.\n\nSee how it works.\n\nBest,\n[Partner Name]"
                }
            ],
            "linkedin_posts": [
                "Stay compliant without the stress. HostBooks automates TDS, GST, and payroll filings dynamically. #TaxCompliance #CorporateGovernance #Enterprise"
            ],
            "whatsapp_messages": [
                "Eliminate tax filing penalties. HostBooks automates statutory tax deadlines: [Link]"
            ],
            "webinar_outline": {
                "title": "Simplifying Corporate Compliance in the Digital Age",
                "agenda": [
                    "Managing shifting TDS/GST tax rates.",
                    "Automating calculations to prevent errors.",
                    "Practice management tools for compliance officers."
                ]
            },
            "outreach_script": "Hi, we help companies streamline their compliance filings. Are you using separate tools for accounting, GST, and TDS filings?"
        }
    ]
    return themes

def get_implementation_success_data():
    return {
        "onboarding": """# Customer Onboarding Playbook

## 1. Phase 1: Initiation & Kickoff
*   **Objective:** Define project team, scope, and key deliverables.
*   **Key Action Items:**
    1.  Hold the kickoff meeting with client and partner leads.
    2.  Set up the HostBooks client tenant and user accounts.
    3.  Confirm migration cutoff date (usually month-end or quarter-end).

## 2. Phase 2: System Configuration
*   **Objective:** Map charts of accounts, cost centers, and user permissions.
*   **Key Action Items:**
    1.  Configure client GSTINs, bank accounts, and invoice templates.
    2.  Define approval matrix roles (Creator, Approver, Audit view).

## 3. Phase 3: Data Migration
*   **Objective:** Import historical master records and balances.
*   *Refer to the Migration Playbook for details.*

## 4. Phase 4: User Training & Dry Run
*   **Objective:** Ensure billing and accounting teams are ready.
*   **Key Action Items:**
    1.  Conduct dry run of GSTR-2B matching and payroll calculations.
    2.  Verify user capability on mobile dashboards.
""",
        "migration": """# Ledger & Compliance Migration Playbook

## 1. Pre-Migration Checklist
*   Export trial balance, vendor ledgers, customer ledgers, and inventory balances from old software.
*   Ensure all bank statement data is complete up to the cutoff date.
*   Review open purchase orders and sales orders.

## 2. Extraction Templates
*   Use the standard HostBooks CSV/Excel templates for:
    *   Chart of Accounts Master
    *   Vendor & Customer Master (with GSTINs, PAN, bank details)
    *   Inventory Item Master (with HSN codes, UOM, tax rates)
    *   Opening Balances Ledger

## 3. Data Ingestion & Validation
*   Run the Migration Wizard tool to import spreadsheets.
*   Validate Trial Balance totals on HostBooks against the source system.
*   *Critical Step:* Cross-check GST ITC opening balance with the government portal dashboard.
""",
        "governance": """# Project Governance Framework

To ensure timely implementations, the following governance model is recommended:

## 1. Roles and Responsibilities
*   **Project Sponsor (Client Executive):** Resolves resource bottlenecks, signs off project milestones.
*   **Implementation Lead (Partner Consultant):** Executes configuration, migration, and training.
*   **Finance Champion (Client Lead Accountant):** Validates balances, manages day-to-day operations.

## 2. Review Cadence
*   **Weekly Progress Sync:** Brief review of task completions, data migrations, and blockers.
*   **Steering Committee Review:** Monthly review of project health and timeline milestones.
""",
        "adoption": """# Customer Adoption Framework

Long-term subscription retention depends on active system adoption.

## 1. Key Adoption Metrics
*   **Login Frequency:** Checking if the finance team logs in daily.
*   **Transaction Volume:** Number of invoices, vouchers, and reconciliations posted.
*   **Feature Utilization:** Are they using automated bank reconciliation and GSP GST filing?

## 2. Playbook for Low Adoption
*   If a client shows low activity within 30 days:
    1. Schedule a system usage check-up meeting.
    2. Re-train team members on automated reconciliation features.
    3. Ensure mobile dashboards are active for owners.
""",
        "qbr_template": """# Quarterly Business Review (QBR) Template

## 1. Executive Summary
*   Reviewing business highlights and financial operational improvements over the past 90 days.

## 2. Realized Value Metrics
*   **Time Reclaimed:** Average bookkeeping close time reduced from X days to Y days.
*   **Compliance Score:** GST filings submitted on time with zero penalties.
*   **ITC Savings:** Value of input tax credit matches recovered by the AI reconciler.

## 3. System Utilization
*   Review feature usage logs (bank feeds, payroll Pay360, e-invoicing).

## 4. Next Quarter Goals
*   Introduce new modules (e.g., adding Payroll to ERP accounts).
*   Optimize inventory replenishment rules.
"""
    }

def get_account_expansion_data():
    return {
        "upsell_cross_sell": """# Upsell & Cross-Sell Playbooks

## 1. Cross-Selling Payroll (Pay360) to Accounting Clients
*   **Target Opportunity:** Accounting clients currently processing payroll in spreadsheets.
*   **Value Pitch:** "Unify employee management with financial tracking. Ledger vouchers for salaries, PF, and TDS post automatically."
*   **Offer:** Free setup and 1-month trial of Pay360 module.

## 2. Upselling Accounting Clients to full ERP
*   **Target Opportunity:** Growing mid-market clients (>50 Cr turnover) experiencing inventory silos and order delays.
*   **Value Pitch:** "Manage procurement, multi-warehouse stock levels, and order fulfillment directly inside your financial software."
""",
        "health_framework": """# Customer Health Framework

Monitor client metrics to predict expansion readiness or churn risk.

| Health Category | Metric Indicator | Status Green | Status Red |
| :--- | :--- | :--- | :--- |
| **System Activity** | User logins per week. | > 4 times per user. | < 1 time per week. |
| **Feature Usage** | Automated bank feeds active. | Active. | Manual imports used. |
| **Compliance Filings** | GST filing status. | Filed on time. | Delayed/Filing errors. |
| **Support Tickets** | Number of unresolved complaints. | 0 to 1 ticket. | > 3 open issues. |
""",
        "triggers": """# Expansion Opportunity Triggers

Look for the following signals in client operations to initiate expansion pitches:

1.  **Multi-Location Expansion:** Client opens a new warehouse or retail branch -> Pitch multi-location inventory and branch accounting.
2.  **Employee Growth:** Headcount crosses 50 employees -> Pitch automated payroll Pay360 integrations.
3.  **Procurement Volume Rise:** Vendor bills increase by 50% -> Pitch AI reconciliation and automated PO matching tools.
4.  **Audit Notices:** Client struggles with statutory compliance review -> Pitch Auditor role access and Compliance suite.
"""
    }

def get_partner_program_data():
    return {
        "onboarding_guide": """# Partner Onboarding Guide

Welcome to the HostBooks Partner Network. This guide outlines your 5-step journey to success.

## 1. Phase 1: Apply
*   Submit your partnership interest form.
*   Review and sign the official partner agreement.

## 2. Phase 2: Onboard
*   Complete KYC verification.
*   Set up your partner account and log in to the Partner Portal.

## 3. Phase 3: Get Enabled
*   Attend structured product training (Functional Consultant, Sales Positioning, and Technical Architecture).
*   Complete assessments to achieve HostBooks Certification.
*   Download sales enablement kits, playbooks, and co-branded marketing campaigns.

## 4. Phase 4: Go Live
*   Access demo environments and development sandbox portals.
*   Receive and register your first client opportunities in the portal.

## 5. Phase 5: Earn & Grow
*   Pitch HostBooks to qualified accounts, close deals, and earn high commissions.
*   Grow your client portfolio to unlock higher partnership tiers and marketing benefits.
""",
        "reseller_referral": """# Partner Models & Scope

Choose the model that fits your organization's strengths, or combine them to maximize growth:

## 1. Solution Partner (Sells & Implements)
*   **Role:** Generate leads, conduct product demonstrations, close sales, manage implementation and client onboarding, and provide first-line support.
*   **Commissions:** **40% Year 1** net contract value + **25% renewals** on subsequent years.
*   **Ideal For:** System Integrators (SIs), CA firms, and IT service providers.

## 2. Sales Partner (Owns Full Sales Cycle)
*   **Role:** Own the sales cycle from lead generation to contract signing. HostBooks handles implementation and onboarding.
*   **Commissions:** **35% Year 1** net contract value + **20% renewals** on subsequent years.
*   **Ideal For:** Sales-led agencies and independent consultants.

## 3. Referral Partner (Introduces Opportunities)
*   **Role:** Introduce qualified leads and coordinate a warm handover to the HostBooks team. HostBooks manages the sale, closing, and implementation.
*   **Commissions:** **10% Year 1** net contract value (Nil renewals).
*   **Ideal For:** Business advisors and industry influencers.
""",
        "deal_reg": """# Deal Registration Process

Secure and protect your sales pipeline using our portal:

1.  **Register Lead:** Log in to the HostBooks Partner Portal, navigate to 'Register Deal', and enter client contact and requirements.
2.  **Duplicate Check:** The system verifies there are no active registrations for that account.
3.  **Fast Approval:** Partner Account Managers review and approve registrations within 4 hours.
4.  **Validity Period:** Registered deals are reserved for the registering partner for 90 days, with options for extension based on progress.
""",
        "incentives_faq": """# Incentives & Partner FAQs

## 1. Partner Enablement Benefits
HostBooks invests directly in your team's capability to close mid-market and enterprise deals:
*   **Commissions Payout:** Up to 40% first-year commissions plus 25% recurring renewals (compounds year-over-year).
*   **Co-Branded Marketing:** Pre-designed GTM campaigns, landing page templates, email templates, and joint event sponsorship funds.
*   **Lead Sharing:** Direct routing of qualified inbound enterprise inquiries to active certified partners.
*   **Demo & Sandbox Access:** High-performance demo instances showcasing 31 modules and 5 industries.
*   **Dedicated Support:** Direct access to a Partner Account Manager (PAM) and pre-sales engineering resources.

## 2. Corporate Office Locations
*   **India:** Nimai Tower, 412-415, Udyog Vihar Phase 4, Gurugram, Haryana 122016
*   **UAE:** Office 29-31, 10th Floor, Ibn Batuta Gate Offices, Jebel Ali, Dubai, United Arab Emirates
*   **USA:** 180 Promenade Cir. Ste 300, Sacramento, CA 95834

## 3. Frequently Asked Questions (FAQ)
*   **Q: How is commission calculated?** It is calculated as a percentage of the net subscription/contract value.
*   **Q: Do you support data migrations from SAP or Tally?** Yes, HostBooks provides automated data migration templates and specialized engineering support.
*   **Q: Can we customize document layouts for clients?** Yes, HostBooks includes highly configurable designers for invoices, BOM forms, and reports.
""",
        "certification": """# Certification Framework

Build expertise and trust with certified qualifications:

## 1. Certified Functional Consultant (ERP-FC)
*   *Focus:* Multi-warehouse SCM configurations, production planning, BOM setup, and multi-entity accounts consolidation.

## 2. Certified Compliance Specialist (ERP-CS)
*   *Focus:* Government GSP connections, AI-powered GSTR-2B matching, TDS deductions, and automated E-Invoicing workflows.
"""
    }

def get_templates_training_case_studies():
    return {
        "templates": """# Standard Templates & Worksheets

1.  [Partner Outreach Script](repository/08%20Templates/partner_outreach_script.docx)
2.  [Client Discovery & Follow-up Questionnaire](repository/08%20Templates/client_discovery_followup.docx)
3.  [Client Requirements Mapping & BRD Guide](repository/08%20Templates/brd_requirements_mapping.docx)
4.  [Customer QBR Review Presentation (PowerPoint)](repository/08%20Templates/qbr_presentation.pptx)
""",
        "training": """# Partner Training Curriculum

## Module 1: System Basics & Navigation (2 Hours)
*   Understanding the user dashboard.
*   Setting up client profiles and user permissions.

## Module 2: Advanced GST & E-Invoicing (3 Hours)
*   Connecting with GSP APIs.
*   Reconciling GSTR-2B automatically.
*   Handling mismatches and vendor notices.

## Module 3: Cloud ERP & Inventory (4 Hours)
*   Configuring bill of materials (BOM).
*   Managing multi-warehouse stock entries.
*   Closing month-end books.

# [Official HostBooks Video Playlists (Click to view all) ↗](https://www.youtube.com/@HostBooksLimited/playlists)
## 1. [HostBooks ERP360 & SCM Tutorials](https://www.youtube.com/@HostBooksLimited/search?query=ERP)
## 2. [Aahar POS & F&B Management Guides](https://www.youtube.com/@HostBooksLimited/search?query=Aahar)
## 3. [Unfiltered with HostBooks (Industry Insights)](https://www.youtube.com/@HostBooksLimited/search?query=Unfiltered)
## 4. [GST, E-Invoicing & Compliance Guides](https://www.youtube.com/@HostBooksLimited/search?query=GST)
""",
        "case_studies": """# [Client Success Case Studies (Click to view all) ↗](https://www.hostbooks.com/in/case-study/)

## 1. [Inde Hotels & Resorts (Hospitality Vertical)](https://www.hostbooks.com/in/case-study/inde-hotels-and-resorts)
*   **Client Overview:** A premium hotel group managing multiple boutique and heritage properties across India.
*   **The Challenge:** Fragmented accounting systems per hotel property led to a labyrinth of complexities; lack of consolidated group financial reports, inventory management discrepancies, and high dependency on frequent physical audits.
*   **Solution Implemented:** Deployed **HostBooks FnB360** to centralize operational control, integrate PMS room/banquet billing with accounts, and automate Input Tax Credit calculations for raw procurement.
*   **The Impact:** Reduced monthly financial closing cycles by **60%**, consolidated balance sheets instantly, and streamlined property audit compliance.

## 2. [Ideal Ice Cream (Manufacturing & SCM)](https://www.hostbooks.com/in/case-study/ice-cream-craftsman)
*   **Client Overview:** Leading ice cream manufacturer and distributor network (established in 1975) operating through over 4,500+ retailers.
*   **The Challenge:** Fragmented applications for factory, dealers, and finance; manual warehouse management led to order fulfillment delays, inaccurate inventory levels, and logistics tracking issues.
*   **Solution Implemented:** Migrated SCM, warehouse management, and production scheduling to **HostBooks ERP360**.
*   **The Impact:** Consolidated distributor orders automatically, leading to a **40% increase in order fulfillment speed**, daily cost monitoring at every operational level, and optimized SKU-wise tracking.

## 3. [Gwalia Sweets (Sweets & Bakeries)](https://www.hostbooks.com/in/case-study/sweets-bakeries)
*   **Client Overview:** Established in 1994, a leading sweets and confectionery brand from Gujarat, managing over 30 retail outlets ranging up to 10,000 square feet.
*   **The Challenge:** Wrestled with multiple systems, manual processes, and lack of a centralized control hub, leading to delayed processes, stock discrepancies, and reporting lags.
*   **Solution Implemented:** Implemented **HostBooks FnB360** with centralized recipe management, integrated warehouse-to-outlet SCM, and automated outlet POS systems.
*   **The Impact:** Achieved an **80% reduction in report-related issues**, established absolute recipe uniformity across outlets, and reduced inventory management time.

## 4. [Sharman Jain Sweets (Traditional Sweets & Confectionery)](https://www.hostbooks.com/in/case-study/sweets-confectioners)
*   **Client Overview:** Distinguished traditional dessert producer and retailer (established in 1995) with owned and franchised outlets across Punjab.
*   **The Challenge:** Highly perishable inventory with high wastage risks, coordination issues between owned and franchised retail outlets, and lack of real-time sales visibility.
*   **Solution Implemented:** Deployed **HostBooks FnB360** with Trend, Festivities, and Sales Analytics, centralized recipe manager, and multi-outlet retail POS.
*   **The Impact:** Standardized quality control procedures, optimized raw material tracking, and minimized ingredient/shelf wastage by **35%**.

## 5. [Automation & Robotic (Special Purpose Machinery)](https://www.hostbooks.com/in/case-study/automation-and-robotic)
*   **Client Overview:** A leading automation engineering firm specializing in Special Purpose Machinery (SPM) serving automotive, biotech, and pharmaceutical industries across Coimbatore, Chennai, Bangalore, and Dubai.
*   **The Challenge:** Lack of project cost transparency, delayed project milestones due to last-minute fire-fighting, manual tracking of materials, and heavy reliance on external tax consultants.
*   **Solution Implemented:** Integrated **HostBooks PMS360 & ERP360** for project timeline tracking, project-based cost center mapping, and compliance automation.
*   **The Impact:** Decreased production times, minimized inventory holding, established a unified data source, and achieved automated GST/TDS reconciliation.

## 6. [Engineering Solution Provider (Automation & Integration)](https://www.hostbooks.com/in/case-study/engineering-solution-provider)
*   **Client Overview:** Established in 2009, a major Haryana-based integration company with three operational plants and upcoming sites in Rajasthan and Gujarat.
*   **The Challenge:** Siloed data across plants, manual project tracking causing delivery delays, manual job worker material reconciliation, and lack of tool room visibility.
*   **Solution Implemented:** Deployed **HostBooks ERP360** to track machine, labor, and tool efficiency; centralized proposal management.
*   **The Impact:** Automated proposal generation from historical data, real-time project-wise inventory tracking, and seamless compliance management with no additional overheads.

## 7. [Luxury Property Construction (Real Estate & Construction)](https://www.hostbooks.com/in/case-study/luxury-property-construction)
*   **Client Overview:** A prominent luxury property developer with over 12 lakh sq. ft. of residential and commercial spaces constructed.
*   **The Challenge:** Inconsistent inventory levels across remote construction sites, difficulties tracking project unit economics, and manual coordination of vendor contracts.
*   **Solution Implemented:** Implemented **HostBooks ERP360 Construction Module** with automated purchasing control and sales forecasting tools.
*   **The Impact:** Enabled precise project-wise cost tracking, optimized material holding costs, and automated subcontractor payment compliance.
"""
    }

def write_repository_files():
    uh = get_understand_hostbooks_data()
    for name, content in uh.items():
        with open(os.path.join(REPO_DIR, "01 Understand HostBooks", f"{name}.md"), "w") as f:
            f.write(content)

    io = get_opportunity_opportunities_data()
    with open(os.path.join(REPO_DIR, "02 Identify Opportunities", "icps.md"), "w") as f:
        f.write(io["icps"])
    with open(os.path.join(REPO_DIR, "02 Identify Opportunities", "opportunity_maps.md"), "w") as f:
        f.write(io["opportunity_maps"])
    
    dq_md = "# Discovery Question Library\n\n"
    current_cat = ""
    for item in io["discovery_questions"]:
        if item["category"] != current_cat:
            current_cat = item["category"]
            dq_md += f"\n## {current_cat} Questions\n\n"
        dq_md += f"*   **{item['id']}**: {item['question']}\n"
    with open(os.path.join(REPO_DIR, "02 Identify Opportunities", "discovery_questions.md"), "w") as f:
        f.write(dq_md)

    sh = get_sell_hostbooks_data()
    with open(os.path.join(REPO_DIR, "03 Sell HostBooks", "sales_playbook.md"), "w") as f:
        f.write(sh["sales_playbook"])
    with open(os.path.join(REPO_DIR, "03 Sell HostBooks", "roi_framework.md"), "w") as f:
        f.write(sh["roi_framework"])
    
    obj_md = "# Objection Handling Library\n\n"
    for item in sh["objections"]:
        obj_md += f"## {item['id']}: {item['objection']}\n"
        obj_md += f"*   **Root Concern:** {item['concern']}\n"
        obj_md += f"*   **Recommended Response:** {item['response']}\n"
        obj_md += f"*   **Supporting Proof Point:** {item['proof']}\n\n"
    with open(os.path.join(REPO_DIR, "03 Sell HostBooks", "objections.md"), "w") as f:
        f.write(obj_md)

    mt = get_marketing_toolkit_data()
    for campaign in mt:
        camp_md = f"# Campaign Kit: {campaign['name']}\n\n"
        camp_md += f"## Campaign Strategy\n{campaign['strategy']}\n\n"
        camp_md += f"## Landing Page Copy\n"
        camp_md += f"*   **Headline:** {campaign['landing_page']['headline']}\n"
        camp_md += f"*   **Subheadline:** {campaign['landing_page']['subheadline']}\n"
        camp_md += f"*   **Benefits:**\n"
        for ben in campaign['landing_page']['benefits']:
            camp_md += f"    * {ben}\n"
        camp_md += f"*   **Call To Action:** {campaign['landing_page']['cta']}\n\n"
        
        camp_md += f"## Email Sequence\n"
        for i, email in enumerate(campaign['email_sequence']):
            camp_md += f"### Email {i+1}: {email['subject']}\n"
            camp_md += f"```text\n{email['body']}\n```\n\n"
        
        camp_md += f"## LinkedIn Posts\n"
        for post in campaign['linkedin_posts']:
            camp_md += f"*   {post}\n"
        camp_md += f"\n"
        
        camp_md += f"## WhatsApp Templates\n"
        for msg in campaign['whatsapp_messages']:
            camp_md += f"*   \"{msg}\"\n"
        camp_md += f"\n"
        
        camp_md += f"## Webinar Outline\n"
        camp_md += f"*   **Title:** {campaign['webinar_outline']['title']}\n"
        camp_md += f"*   **Agenda:**\n"
        for ag in campaign['webinar_outline']['agenda']:
            camp_md += f"    * {ag}\n"
        camp_md += f"\n"
        
        camp_md += f"## Partner Outreach Script\n"
        camp_md += f"```text\n{campaign['outreach_script']}\n```\n"

        with open(os.path.join(REPO_DIR, "04 Marketing Toolkit", f"{campaign['id']}.md"), "w") as f:
            f.write(camp_md)

    im = get_implementation_success_data()
    for name, content in im.items():
        with open(os.path.join(REPO_DIR, "05 Implementation Success", f"{name}.md"), "w") as f:
            f.write(content)

    ae = get_account_expansion_data()
    for name, content in ae.items():
        with open(os.path.join(REPO_DIR, "06 Account Expansion", f"{name}.md"), "w") as f:
            f.write(content)

    pp = get_partner_program_data()
    for name, content in pp.items():
        with open(os.path.join(REPO_DIR, "07 Partner Program", f"{name}.md"), "w") as f:
            f.write(content)

    oth = get_templates_training_case_studies()
    with open(os.path.join(REPO_DIR, "08 Templates", "templates_index.md"), "w") as f:
        f.write(oth["templates"])

    # Write physical high-fidelity template files so they can be served and downloaded
    outreach_md = """# HostBooks Partner Outreach Script
This script is designed for cold and warm outreach to CFOs, IT Directors, and Finance Heads of mid-market and enterprise accounts (50 Cr+ turnover).

---

## 1. Cold Email Template: ERP Modernization & TCO Reduction
**Subject:** 35% Lower TCO for your ERP & Compliance: HostBooks Autonomous Agentic ERP

Dear [CFO_Name],

I hope this email finds you well.

I am reaching out from [Partner_Company], an enterprise solutions partner. We assist mid-market organizations in transitioning from high-cost, fragmented legacy software to high-efficiency, automated cloud systems.

Many finance heads of companies with a turnover exceeding 50 Cr face similar operational challenges:
*   **Manual Reconciliations:** Teams spending days matching purchase invoices with GSTR-2B.
*   **Siloed Systems:** Separate payroll (HRMS), warehouse inventory, and tax compliance modules.
*   **High Licensing Fees:** Inflated annual renewal fees for SAP, Oracle, or Microsoft Dynamics.

**HostBooks is the world's first Autonomous Agentic ERP**, running 24x7 autonomous AI agents that handle multi-way AP matching, automated inventory stock replenishment, and direct GST compliance. On average, our clients reduce total cost of ownership (TCO) by **35% to 55%** while going live in a matter of weeks, not months.

Are you available for a brief 10-minute introduction call next Tuesday at 11:00 AM?

Best regards,

[Partner_Name]  
[Partner_Title]  
[Partner_Company]  
[Contact_Number]

---

## 2. LinkedIn Outreach Message (Short & Impactful)
"Hello [First_Name], noticed your work leading finance operations at [Company]. Many CFOs of 50 Cr+ businesses tell us they are looking to trim legacy ERP license overheads (SAP/Dynamics) and automate AP matching. HostBooks runs 24x7 AI agents that cut ERP license costs by 40% and automate GST compliance. I'd love to share our 1-page comparison sheet. Do you have a brief moment to connect next week?"

---

## 3. Cold Calling Script & Elevator Pitch
*   **Hook:** "Hello [CFO_Name], my name is [Your_Name] from [Partner_Company]. I'm calling because we help mid-market companies in [Client_Industry] automate their finance operations and cut legacy software costs by 35%."
*   **Value Proposition:** "We implement HostBooks—the world's first Autonomous Agentic ERP. Instead of manual data entry, it uses background AI agents to reconcile bank statements and GSTR-2B automatically in real time."
*   **Qualification Question:** "Are you currently running on legacy systems like Tally or SAP, and how many hours does your team spend weekly on GST matching?"
*   **Call-to-Action:** "I'd like to schedule a 15-minute live demo showing how we automate these workflows for businesses in your turnover range. Would Thursday afternoon work?"
"""
    save_as_docx("HostBooks Partner Outreach Script", outreach_md, os.path.join(REPO_DIR, "08 Templates", "partner_outreach_script.docx"))

    discovery_md = """# Client Discovery & Follow-up Questionnaire
Use this worksheet to gather critical technical and business details about the client's internal operations and workflows.

---

## 1. Corporate Profile & System Landscape
*   **Current Primary Systems:** [ ] Tally  [ ] SAP  [ ] Zoho  [ ] MS Dynamics  [ ] In-house Custom
*   **Hosting Model:** [ ] On-Premise  [ ] Private Cloud  [ ] SaaS / Multi-tenant Cloud
*   **Turnover Range:** [ ] 50 Cr - 100 Cr  [ ] 100 Cr - 250 Cr  [ ] 250 Cr - 500 Cr+
*   **Number of Legal Entities / Branches:** ________ / ________
*   **Key Integrations Needed:** [ ] E-Commerce  [ ] CRM  [ ] Third-party Logistics (3PL)  [ ] Banks

---

## 2. Finance & Tax Compliance Gaps
1.  **GSTR-2B Matching Frequency:**
    *   [ ] Daily  [ ] Weekly  [ ] Monthly at closing  [ ] Quarterly
2.  **Purchase Invoice Volume:** ____________ invoices per month.
3.  **Manual Labor in Accounts Payable (AP):**
    *   How many employees are dedicated solely to manual invoice entry, validation, and multi-way matching? ________ FTEs.
4.  **GST Reconciliation Pain Points:**
    *   Are you losing Input Tax Credit (ITC) due to vendor non-compliance? [ ] Yes  [ ] No
    *   Estimated annual ITC leakage: ____________ INR.

---

## 3. Supply Chain & Manufacturing Gaps
1.  **Warehouse Management:**
    *   Number of warehouses: ________
    *   Are stock updates real-time? [ ] Yes  [ ] No (Manual batch uploads)
2.  **Inventory Tracking:**
    *   Do you experience SKU discrepancies during audits? [ ] Frequently  [ ] Occasionally  [ ] Never
    *   Do you support multi-level Bill of Materials (BOM)? [ ] Yes  [ ] No
3.  **Purchase Approval Workflows:**
    *   Is purchase requisition, approval, and purchase order creation automated? [ ] Yes  [ ] No

---

## 4. HR, Payroll & Operations
1.  **Total Employee Count:** ________
2.  **Payroll Calculation Method:**
    *   [ ] Automated inside main system  [ ] Independent HRMS  [ ] Excel spreadsheets
3.  **Biometric & Attendance Integration:**
    *   Do you require biometric machine sync for payroll deductions? [ ] Yes  [ ] No
"""
    save_as_docx("Client Discovery & Follow-up Questionnaire", discovery_md, os.path.join(REPO_DIR, "08 Templates", "client_discovery_followup.docx"))

    brd_md = """# Client Requirements Mapping & BRD Guide
This guide outlines how to translate the client's operational pain points into a formal Business Requirements Document (BRD) and map them to HostBooks ERP modules.

---

## 1. The Requirements Mapping Matrix
Use this framework to map the client's internal operational challenges directly to HostBooks features.

| Identified Pain Point | Business Impact | HostBooks ERP Module | Desired Outcome |
| :--- | :--- | :--- | :--- |
| **Manual AP Entry & 2B Matching** | ITC leakage, human typing errors, delayed monthly books closing. | **Autonomous AP Agents & GSP Compliance** | 99.8% reconciliation accuracy, instant ITC recovery, 60% faster closing. |
| **Warehouse stock discrepancies** | Production delays, overstocking capital lockup, delivery lag. | **HostBooks SCM & Multi-Warehouse Control** | Real-time SKU tracking, automated reorder triggers, centralized replenishment. |
| **Delayed management dashboards** | Delayed strategic decisions, disconnected sales/finance reports. | **Centralized BI & Single Ledger Database** | Real-time consolidated balance sheets and SKU profitability margins. |
| **Manual payroll calculations** | Attendance discrepancies, compliance delays (PF/ESIC), human calculation errors. | **HostBooks Pay360 (HR & Payroll)** | Biometric attendance sync, automated deductions, direct salary bank transfers. |

---

## 2. Structure of a Business Requirements Document (BRD)
When drafting the BRD for a 50 Cr+ client, structure the document as follows:

### Section 1: Executive Summary & Project Scope
*   Brief overview of the client's business model (e.g. F&B retail chain, specialized manufacturing plant).
*   Project goals (e.g. migrate 3 legal entities from on-premise Tally to HostBooks cloud ERP).

### Section 2: Current State vs Future State (Process Maps)
*   **AS-IS Workflow:** Show manual steps, email dependencies, and Excel spreadsheets.
*   **TO-BE Workflow:** Show direct HostBooks modules, automated background AI agent operations, and integrated real-time compliance.

### Section 3: Detailed Functional Requirements
1.  **General Ledger & Consolidation:** Inter-company transactions, multi-currency reporting.
2.  **Inventory & Procurement:** Purchase approvals hierarchy, automated GRN (Goods Received Note).
3.  **Sales & Retail POS:** In-store POS syncing, central price-list pushing.
4.  **Taxation & Compliance:** Automatic generation of e-invoices and e-way bills.

### Section 4: Data Migration & Implementation Plan
*   Details on importing ledger masters, customer database, vendor details, and opening inventory balances.
*   UAT (User Acceptance Testing) guidelines and training schedules.
"""
    save_as_docx("Client Requirements Mapping & BRD Guide", brd_md, os.path.join(REPO_DIR, "08 Templates", "brd_requirements_mapping.docx"))

    qbr_ppt = """Slide 1: Title
- Client Quarterly Business Review (QBR)
- Date: Q3 / Q4 Evaluation
- Presented by: [Partner Name]
- Product: HostBooks Enterprise ERP & Compliance Suite

Slide 2: Agenda
- Executive Summary & Success Stories
- Financial Performance & GST Compliance Metrics
- Operational Efficiency: SCM, Warehouse & Recipe Controls
- Account Expansion & Future Roadmap

Slide 3: Key Performance Metrics
- Financial Closing Cycle reduced by 60%
- GSTR Reconciliation automation accuracy at 99.8%
- Food Recipe Control wastage reduced by 35%
- Stock Replenishment efficiency boosted by 40%

Slide 4: Upcoming Milestones
- Phase 2: FnB POS Integration across all Franchise outlets
- Phase 3: Project-based Cost Center tracking inside PMS360
- CRM Lead generation integrations

Slide 5: Open Discussion & Action Plan
- Q&A session
- Next review date scheduled
- Support & escalations matrix
"""
    save_as_pptx("Client Quarterly Business Review (QBR)", qbr_ppt, os.path.join(REPO_DIR, "08 Templates", "qbr_presentation.pptx"))

    with open(os.path.join(REPO_DIR, "09 Training", "training_curriculum.md"), "w") as f:
        f.write(oth["training"])
    with open(os.path.join(REPO_DIR, "10 Case Studies", "case_studies.md"), "w") as f:
        f.write(oth["case_studies"])

    print("Created all markdown repository files.")

def generate_database_js():
    io = get_opportunity_opportunities_data()
    sh = get_sell_hostbooks_data()
    
    data = {
        "understand": get_understand_hostbooks_data(),
        "opportunities": {
            "icps": io["icps"],
            "opportunity_maps": io["opportunity_maps"],
            "discovery_questions": io["discovery_questions"]
        },
        "sell": {
            "sales_playbook": sh["sales_playbook"],
            "objections": sh["objections"],
            "roi_framework": sh["roi_framework"]
        },
        "marketing": get_marketing_toolkit_data(),
        "implementation": get_implementation_success_data(),
        "expansion": get_account_expansion_data(),
        "partner_program": get_partner_program_data(),
        "templates": {
            "templates_index": get_templates_training_case_studies()["templates"]
        },
        "training": {
            "training_curriculum": get_templates_training_case_studies()["training"]
        },
        "case_studies": {
            "case_studies": get_templates_training_case_studies()["case_studies"]
        }
    }
    
    js_content = f"// Automatically generated partner database\nconst PARTNER_DATA = {json.dumps(data, indent=2)};\n"
    with open(os.path.join(BASE_DIR, "data.js"), "w") as f:
        f.write(js_content)
    print("Generated data.js database file.")

if __name__ == "__main__":
    create_directory_structure()
    write_repository_files()
    generate_database_js()
